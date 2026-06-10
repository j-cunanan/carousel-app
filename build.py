#!/usr/bin/env python3.10
"""IG carousel build pipeline.

    carousel_src.html + assets/archivo.css
        -> carousel.html          (fonts inlined, deterministic render)
        -> out/slide_NN.png       (1080x1350 each, via Playwright)
        -> out/carousel.pptx      (one slide per page, for Canva import)

Usage:
    uv run python build.py            # full build
    uv run python build.py --scale 2  # 2160x2700 retina PNGs (PPTX stays 1080x1350)
    uv run python build.py --video-source <url-or-file>  # also export branded MP4

Asset capture (run once per story):
    uv run python capture_tweet_page.py <tweet_url> assets/tweet_original.png
"""
import argparse
import sys
from pathlib import Path

ROOT = Path(__file__).parent
SRC = ROOT / "carousel_src.html"
FONTS = ROOT / "assets" / "archivo.css"
HTML = ROOT / "carousel.html"
OUT = ROOT / "out"

SLIDE_W, SLIDE_H = 1080, 1350
EMU_PER_PX = 9525  # 96 dpi


def inject_fonts() -> None:
    html = SRC.read_text().replace("/* __FONTS__ */", FONTS.read_text())
    HTML.write_text(html)
    print(f"[1/3] fonts inlined -> {HTML.name}")


def render_slides(scale: int) -> list[Path]:
    from playwright.sync_api import sync_playwright

    OUT.mkdir(exist_ok=True)
    pngs = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": SLIDE_W + 100, "height": SLIDE_H + 100},
            device_scale_factor=scale,
        )
        page.goto(HTML.as_uri())
        page.wait_for_load_state("networkidle")
        page.wait_for_timeout(400)  # font/image settle

        slides = page.locator(".slide")
        count = slides.count()
        for i in range(count):
            out_path = OUT / f"slide_{i + 1:02d}.png"
            slides.nth(i).screenshot(path=str(out_path))
            pngs.append(out_path)
            print(f"[2/3] rendered {out_path.name} ({SLIDE_W * scale}x{SLIDE_H * scale})")
        browser.close()
    return pngs


def export_pptx(pngs: list[Path]) -> Path:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * EMU_PER_PX)
    prs.slide_height = Emu(SLIDE_H * EMU_PER_PX)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    out_path = OUT / "carousel.pptx"
    prs.save(out_path)
    print(f"[3/3] exported {out_path} ({len(pngs)} pages)")
    return out_path


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--scale", type=int, default=1, help="PNG pixel density multiplier")
    ap.add_argument("--video-source", help="Also export a branded MP4 slide from this local file or URL")
    ap.add_argument("--video-tweet-embed-file", type=Path, help="HTML file containing a twitter-tweet blockquote")
    ap.add_argument("--video-layout", choices=["video", "post-video"], default="video")
    ap.add_argument("--video-out", type=Path, default=OUT / "video_slide_02.mp4")
    ap.add_argument("--video-frame-out", type=Path, default=OUT / "video_frame_02.png")
    ap.add_argument("--video-poster-out", type=Path, default=OUT / "video_slide_02_poster.png")
    ap.add_argument("--video-caption", default="A source video, framed inside the LLMAW carousel system.")
    ap.add_argument("--video-kicker", default="Video receipt")
    ap.add_argument("--video-source-label", default="")
    ap.add_argument("--video-active", type=int, default=2, help="Active dot index for the MP4 slide")
    ap.add_argument("--video-count", type=int, default=6, help="Total progress dots for the MP4 slide")
    ap.add_argument("--video-fit", choices=["contain", "cover"], default="contain")
    ap.add_argument("--video-fps", type=int, default=30)
    ap.add_argument("--video-mute", action="store_true")
    ap.add_argument(
        "--video-cookies-from-browser",
        help="Pass through to yt-dlp, for example chrome or safari when X gates media",
    )
    ap.add_argument("--video-post-author", help="Override detected post author in post-video layout")
    ap.add_argument("--video-post-handle", help="Override detected post handle in post-video layout")
    ap.add_argument("--video-post-text", help="Override detected post text in post-video layout")
    ap.add_argument("--video-post-date", help="Override detected post date in post-video layout")
    args = ap.parse_args()

    inject_fonts()
    pngs = render_slides(args.scale)
    if not pngs:
        print("no .slide elements found", file=sys.stderr)
        return 1
    export_pptx(pngs)

    if args.video_source or args.video_tweet_embed_file:
        from build_video_slide import build_video_slide

        build_video_slide(
            source=args.video_source,
            tweet_embed_file=args.video_tweet_embed_file,
            out_path=args.video_out,
            frame_out=args.video_frame_out,
            poster_out=args.video_poster_out,
            caption=args.video_caption,
            kicker=args.video_kicker,
            source_label=args.video_source_label,
            active=args.video_active,
            count=args.video_count,
            fit=args.video_fit,
            fps=args.video_fps,
            mute=args.video_mute,
            cookies_from_browser=args.video_cookies_from_browser,
            layout=args.video_layout,
            post_author=args.video_post_author,
            post_handle=args.video_post_handle,
            post_text=args.video_post_text,
            post_date=args.video_post_date,
        )
    return 0


if __name__ == "__main__":
    sys.exit(main())
